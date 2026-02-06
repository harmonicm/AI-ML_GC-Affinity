from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


def render_ppt(teaser_json, output_path):
    prs = Presentation()

    for slide_data in teaser_json["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data["title"]

        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        for i, bullet in enumerate(slide_data["bullets"]):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = bullet
            p.level = 0

        if slide_data.get("chart"):
            chart = slide_data["chart"]
            data = ChartData()
            data.categories = chart["labels"]
            data.add_series("Series", chart["values"])

            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(4), Inches(1.5),
                Inches(4), Inches(3),
                data
            )

    prs.save(output_path)
