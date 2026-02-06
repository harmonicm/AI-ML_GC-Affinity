import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


# =========================
# KELP BRAND SYSTEM
# =========================
KELP_DARK = RGBColor(43, 36, 79)
TEXT_GREY = RGBColor(64, 64, 64)
FOOTER_TEXT = "Strictly Private & Confidential – Prepared by Kelp M&A Team"


# =========================
# HELPERS
# =========================
def anonymize(text: str) -> str:
    if not text:
        return ""
    for name in ["Kalyani", "Gati", "Ksolves", "Centum", "Connplex", "Ind Swift"]:
        text = text.replace(name, "The Company")
    return text.strip()


def extract_timeseries(sections, metric_name):
    """
    Extracts last 3 years of numeric data from Financials Status
    """
    financials = sections.get("Financials Status", "")
    pattern = rf"{metric_name}.*"
    match = re.search(pattern, financials)
    if not match:
        return None

    line = match.group(0)
    pairs = re.findall(r"(\d{4}):\s*([-\d.]+)", line)

    clean = [(y, float(v)) for y, v in pairs if v not in ("None", "nan")]
    if len(clean) < 3:
        return None

    last = clean[-3:]
    years = [y for y, _ in last]
    values = [v for _, v in last]
    return years, values


# =========================
# SLIDE ELEMENTS
# =========================
def add_logo(slide):
    box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(1.5), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = "KELP"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.name = "Arial"
    p.font.color.rgb = KELP_DARK


def add_footer(slide):
    box = slide.shapes.add_textbox(Inches(2.5), Inches(6.9), Inches(5), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = FOOTER_TEXT
    p.font.size = Pt(9)
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1), Inches(8), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.name = "Arial"
    p.font.color.rgb = KELP_DARK


def add_bullets(slide, bullets, x, y, width=3.8):
    bullets = [b.strip()[:110] for b in bullets if b.strip()]
    if not bullets:
        return

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(3))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, b in enumerate(bullets[:4]):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = b
        p.font.size = Pt(11)
        p.font.name = "Arial"
        p.font.color.rgb = TEXT_GREY
        p.space_after = Pt(10)
        p.level = 1


# =========================
# FINAL PPT GENERATOR
# =========================
def generate_ppt_from_sections(sections, output_path, sector=None):
    prs = Presentation()

    # ---------- SLIDE 1 ----------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    add_title(slide, "Company Overview")
    add_bullets(
        slide,
        [
            anonymize(sections.get("Business Description", ""))[:120],
            "Operates a scalable and quality-driven manufacturing platform",
            "Strong presence across domestic and export markets",
        ],
        0.7,
        2,
    )
    add_footer(slide)

    # ---------- SLIDE 2 ----------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    add_title(slide, "Products & End Markets")
    add_bullets(slide, sections.get("Product & Services", "").split("\n"), 0.7, 2)
    add_bullets(
        slide,
        sections.get("Application areas / Industries served", "").split(","),
        4.8,
        2,
    )
    add_footer(slide)

    # ---------- SLIDE 3 (REAL FINANCIAL CHART) ----------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    add_title(slide, "Financial Performance")

    metric_map = {
        "automotive": "Revenue From Operations",
        "technology": "Revenue From Operations",
        "electronics": "Revenue From Operations",
        "logistics": "Revenue From Operations",
        "pharma": "Operating EBITDA",
    }

    metric = metric_map.get(sector, "Revenue From Operations")
    series = extract_timeseries(sections, metric)

    if series:
        years, values = series
        chart_data = ChartData()
        chart_data.categories = years
        chart_data.add_series(metric, values)

        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.6),
            Inches(2.1),
            Inches(5.2),
            Inches(3.1),
            chart_data,
        )

        add_bullets(
            slide,
            [
                f"{metric} shows steady trend over last 3 years",
                "Operational leverage improving",
                "Disciplined cost management",
            ],
            6,
            2,
            2.4,
        )

    add_footer(slide)

    # ---------- SLIDE 4 ----------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    add_title(slide, "Market Opportunity")
    add_bullets(
        slide,
        sections.get("Market Size", "").split("\n")[:4]
        or ["Large and expanding addressable market"],
        0.7,
        2,
    )
    add_footer(slide)

    # ---------- SLIDE 5 ----------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    add_title(slide, "Investment Highlights")
    add_bullets(
        slide,
        [
            "Strong manufacturing and execution capabilities",
            "Diversified customer base with marquee OEMs",
            "Visible growth pipeline and export opportunity",
        ],
        0.7,
        2,
    )
    add_bullets(
        slide,
        [
            "Robust governance and compliance framework",
            "Continuous capacity expansion",
            "Attractive long-term demand drivers",
        ],
        4.8,
        2,
    )
    add_footer(slide)

    prs.save(output_path)
